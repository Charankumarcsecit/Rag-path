"""Document processing utilities for chunking and embedding."""

import os
import logging
import io
from pathlib import Path
from typing import List, Dict, Any, BinaryIO
from dataclasses import dataclass

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


@dataclass
class Document:
    """Represents a document chunk with metadata."""
    content: str
    source: str
    chunk_id: int
    metadata: Dict[str, Any]


class DocumentProcessor:
    """Processes documents for the RAG system."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: Maximum size of each chunk
            chunk_overlap: Overlap between consecutive chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_extensions = {'.md', '.txt', '.py', '.js', '.json', '.yaml', '.yml', 
                                     '.pdf', '.docx', '.pptx', '.html', '.csv'}
    
    def load_documents_from_directory(self, directory: Path) -> List[Document]:
        """
        Load all supported documents from a directory.
        
        Args:
            directory: Path to the directory containing documents
            
        Returns:
            List of Document objects
        """
        documents = []
        
        if not directory.exists():
            logger.warning(f"Directory does not exist: {directory}")
            return documents
        
        for file_path in directory.rglob('*'):
            if file_path.is_file() and file_path.suffix in self.supported_extensions:
                try:
                    docs = self.load_document(file_path)
                    documents.extend(docs)
                except Exception as e:
                    logger.error(f"Error loading {file_path}: {e}")
        
        logger.info(f"Loaded {len(documents)} document chunks from {directory}")
        return documents
    
    def load_document(self, file_path: Path) -> List[Document]:
        """
        Load a single document and split it into chunks.
        
        Args:
            file_path: Path to the document
            
        Returns:
            List of Document chunks
        """
        try:
            # Extract content based on file type
            content = self._extract_content(file_path)
            
            if not content:
                logger.warning(f"No content extracted from {file_path}")
                return []
            
            chunks = self.split_text(content)
            
            documents = []
            for idx, chunk in enumerate(chunks):
                doc = Document(
                    content=chunk,
                    source=str(file_path),
                    chunk_id=idx,
                    metadata={
                        'file_name': file_path.name,
                        'file_type': file_path.suffix,
                        'chunk_index': idx,
                        'total_chunks': len(chunks)
                    }
                )
                documents.append(doc)
            
            return documents
            
        except Exception as e:
            logger.error(f"Error loading document {file_path}: {e}")
            return []
    
    def load_uploaded_file(self, uploaded_file, filename: str) -> List[Document]:
        """
        Load a document from Streamlit uploaded file.
        
        Args:
            uploaded_file: Streamlit UploadedFile object
            filename: Name of the file
            
        Returns:
            List of Document chunks
        """
        try:
            file_ext = Path(filename).suffix.lower()
            content = self._extract_content_from_bytes(uploaded_file.read(), file_ext, filename)
            
            if not content:
                logger.warning(f"No content extracted from {filename}")
                return []
            
            chunks = self.split_text(content)
            
            documents = []
            for idx, chunk in enumerate(chunks):
                doc = Document(
                    content=chunk,
                    source=filename,
                    chunk_id=idx,
                    metadata={
                        'file_name': filename,
                        'file_type': file_ext,
                        'chunk_index': idx,
                        'total_chunks': len(chunks),
                        'upload_type': 'direct_upload'
                    }
                )
                documents.append(doc)
            
            return documents
            
        except Exception as e:
            logger.error(f"Error loading uploaded file {filename}: {e}")
            return []
    
    def _extract_content(self, file_path: Path) -> str:
        """Extract text content from various file types."""
        file_ext = file_path.suffix.lower()
        
        try:
            if file_ext == '.pdf':
                return self._extract_pdf(file_path)
            elif file_ext == '.docx':
                return self._extract_docx(file_path)
            elif file_ext == '.pptx':
                return self._extract_pptx(file_path)
            else:
                # Text-based files
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
        except Exception as e:
            logger.error(f"Error extracting content from {file_path}: {e}")
            return ""
    
    def _extract_content_from_bytes(self, file_bytes: bytes, file_ext: str, filename: str) -> str:
        """Extract text content from file bytes."""
        try:
            if file_ext == '.pdf':
                return self._extract_pdf_from_bytes(file_bytes)
            elif file_ext == '.docx':
                return self._extract_docx_from_bytes(file_bytes)
            elif file_ext == '.pptx':
                return self._extract_pptx_from_bytes(file_bytes)
            else:
                # Text-based files
                return file_bytes.decode('utf-8')
        except Exception as e:
            logger.error(f"Error extracting content from bytes {filename}: {e}")
            return ""
    
    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file."""
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except ImportError:
            logger.warning("pdfplumber not installed, falling back to PyPDF2")
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(str(file_path))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            except Exception as e:
                logger.error(f"Error extracting PDF with PyPDF2: {e}")
                return ""
    
    def _extract_pdf_from_bytes(self, file_bytes: bytes) -> str:
        """Extract text from PDF bytes."""
        try:
            import pdfplumber
            text = ""
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(file_bytes))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            except Exception as e:
                logger.error(f"Error extracting PDF: {e}")
                return ""
    
    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file."""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(str(file_path))
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}")
            return ""
    
    def _extract_docx_from_bytes(self, file_bytes: bytes) -> str:
        """Extract text from DOCX bytes."""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(io.BytesIO(file_bytes))
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}")
            return ""
    
    def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX file."""
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            return ""
    
    def _extract_pptx_from_bytes(self, file_bytes: bytes) -> str:
        """Extract text from PPTX bytes."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            return ""
    
    def split_text(self, text: str) -> List[str]:
        """
        Split text into chunks with overlap.
        
        Args:
            text: Text to split
            
        Returns:
            List of text chunks
        """
        if len(text) <= self.chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + self.chunk_size
            
            # Find the last newline or space before chunk_size
            if end < len(text):
                # Try to break at paragraph
                last_double_newline = text.rfind('\n\n', start, end)
                if last_double_newline > start:
                    end = last_double_newline
                else:
                    # Try to break at sentence
                    last_period = text.rfind('. ', start, end)
                    if last_period > start:
                        end = last_period + 1
                    else:
                        # Try to break at word
                        last_space = text.rfind(' ', start, end)
                        if last_space > start:
                            end = last_space
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - self.chunk_overlap if end < len(text) else end
        
        return chunks
    
    def filter_changed_documents(self, 
                                  all_docs: List[Document], 
                                  changed_files: List[str]) -> List[Document]:
        """
        Filter documents to only those from changed files.
        
        Args:
            all_docs: All documents
            changed_files: List of changed file paths
            
        Returns:
            Filtered list of documents
        """
        changed_docs = []
        
        for doc in all_docs:
            for changed_file in changed_files:
                if changed_file in doc.source:
                    changed_docs.append(doc)
                    break
        
        return changed_docs


def main():
    """Test the document processor."""
    processor = DocumentProcessor()
    
    # Test with the current directory
    docs = processor.load_documents_from_directory(Path('.'))
    
    print(f"Loaded {len(docs)} document chunks")
    if docs:
        print(f"\nFirst chunk:")
        print(f"Source: {docs[0].source}")
        print(f"Content preview: {docs[0].content[:200]}...")


if __name__ == "__main__":
    main()
